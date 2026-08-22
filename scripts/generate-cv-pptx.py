"""Generate a CV-ready PowerPoint for the Flowdesk project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUT = Path(__file__).resolve().parents[1] / "Flowdesk-CV-Portfolio.pptx"

# Brand palette (teal / warm stone — matches product, not purple AI cliché)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x13, 0x4E, 0x4A)
INK = RGBColor(0x1C, 0x19, 0x17)
MUTED = RGBColor(0x57, 0x53, 0x4E)
CREAM = RGBColor(0xFA, 0xF6, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM = RGBColor(0xC2, 0x41, 0x0C)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)


def set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    fill = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return fill


def add_accent_bar(slide, left=0, top=0, width=0.18, height=7.5, color=TEAL):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_textbox(slide, left, top, width, height, paragraphs):
    """paragraphs: list of (text, size, bold, color, align)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, align) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_bullet_block(slide, left, top, width, height, title, bullets, title_color=TEAL_DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, title, size=16, bold=True, color=title_color)
    for b in bullets:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(6)
        r = p.add_run()
        set_run(r, "•  " + b, size=13, color=INK)
    return box


def card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE4, 0xDD, 0xD2)
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ——— 1. Title ———
    s = prs.slides.add_slide(blank)
    add_bg(s, TEAL_DARK)
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.9), Inches(13.333), Inches(1.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_textbox(
        s, 0.9, 1.8, 11.5, 1.2,
        [("FLOWDESK", 48, True, WHITE, PP_ALIGN.LEFT)],
    )
    add_textbox(
        s, 0.9, 2.9, 11.5, 1.0,
        [
            (
                "Family CRM — tickets, calendar, finance, meds, notes & more",
                22,
                False,
                RGBColor(0xCC, 0xFB, 0xF1),
                PP_ALIGN.LEFT,
            )
        ],
    )
    add_textbox(
        s, 0.9, 4.2, 11.5, 0.8,
        [
            (
                "Full-stack product portfolio  ·  Built end-to-end by Arnob Banik",
                16,
                False,
                RGBColor(0xA8, 0xA2, 0x9E),
                PP_ALIGN.LEFT,
            )
        ],
    )
    add_textbox(
        s, 0.9, 6.25, 11.5, 0.8,
        [
            (
                "Live: https://flowdesk-banik.vercel.app   ·   Next.js · TypeScript · Prisma · Neon · Capacitor",
                14,
                False,
                WHITE,
                PP_ALIGN.LEFT,
            )
        ],
    )

    # ——— 2. Project snapshot ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.6,
        [("Project snapshot", 32, True, INK, PP_ALIGN.LEFT)],
    )
    add_textbox(
        s, 0.6, 0.95, 12, 0.5,
        [
            (
                "A production household CRM used for real family planning — shipped as web PWA + Android APK.",
                15,
                False,
                MUTED,
                PP_ALIGN.LEFT,
            )
        ],
    )

    snapshots = [
        ("Problem", "Families juggle tickets, events, meds, money, and notes across chats and apps."),
        ("Solution", "One shared workspace: Today radar, board, calendar, finance, meds, poems, and alerts."),
        ("Scope", "Auth, multi-user household, 40+ API routes, mobile shell, theming, offline-capable PWA."),
        ("Outcome", "Deployed on Vercel + Neon Postgres; Android Capacitor app loads live production URL."),
    ]
    for i, (title, body) in enumerate(snapshots):
        col = i % 2
        row = i // 2
        left = 0.55 + col * 6.3
        top = 1.7 + row * 2.5
        card(s, left, top, 6.0, 2.2)
        add_textbox(
            s, left + 0.35, top + 0.35, 5.3, 0.45,
            [(title, 18, True, TEAL, PP_ALIGN.LEFT)],
        )
        add_textbox(
            s, left + 0.35, top + 0.9, 5.3, 1.0,
            [(body, 14, False, INK, PP_ALIGN.LEFT)],
        )

    # ——— 3. Skills ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.6,
        [("Skills demonstrated", 32, True, INK, PP_ALIGN.LEFT)],
    )

    skill_cols = [
        (
            "Frontend",
            [
                "Next.js 15 (App Router)",
                "React 19 + TypeScript",
                "Tailwind CSS v4",
                "Framer Motion UI",
                "Tiptap rich text",
                "dnd-kit drag & drop",
                "Responsive + PWA (Serwist)",
            ],
        ),
        (
            "Backend & data",
            [
                "REST API route handlers",
                "Prisma ORM + migrations",
                "Neon serverless Postgres",
                "Zod validation",
                "Auth.js (credentials + OAuth)",
                "Cron / alert dispatch",
                "Web Push + email alerts",
            ],
        ),
        (
            "Product & shipping",
            [
                "End-to-end feature design",
                "Household multi-tenancy",
                "Capacitor Android APK",
                "Vercel production deploys",
                "Geolocation / weather APIs",
                "Accessibility & contrast",
                "Command palette (⌘K)",
            ],
        ),
    ]
    for i, (title, items) in enumerate(skill_cols):
        left = 0.55 + i * 4.2
        card(s, left, 1.2, 4.0, 5.6)
        add_bullet_block(s, left + 0.3, 1.45, 3.4, 5.1, title, items)

    # ——— 4. Product modules ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.5,
        [("What I built — product modules", 30, True, INK, PP_ALIGN.LEFT)],
    )

    modules = [
        ("Today Radar", "Due work, meds, events, smart composer, weather, daily quote"),
        ("Tickets & Board", "Kanban statuses, assignees, focus, inbox triage, templates"),
        ("Calendar", "Events, postpone / finish-early, Google sync hooks, holidays"),
        ("Finance", "Income/spend, budgets, savings goals, BD income-tax calculator"),
        ("Meds", "Schedules, dose logs, timezone-aware reminders before dose"),
        ("Notes & Poems", "TipTap notes; poems atelier with dictionary + doodle pad"),
        ("People & Projects", "Contacts, follow-ups, life-area projects"),
        ("Focus", "Pomodoro timer with ticket linkage and stats"),
        ("Alerts", "In-app, Web Push, and email for due items"),
        ("Theme system", "Palettes, light/dark, Bing wallpaper live/pinned modes"),
        ("Auth & household", "Signup, invite codes, Google OAuth, account settings"),
        ("Mobile", "Capacitor Android shell + geolocation + notifications"),
    ]
    for i, (t, d) in enumerate(modules):
        col = i % 3
        row = i // 3
        left = 0.5 + col * 4.2
        top = 1.1 + row * 1.45
        card(s, left, top, 4.0, 1.3)
        add_textbox(
            s, left + 0.25, top + 0.2, 3.5, 0.35,
            [(t, 14, True, TEAL_DARK, PP_ALIGN.LEFT)],
        )
        add_textbox(
            s, left + 0.25, top + 0.55, 3.5, 0.6,
            [(d, 12, False, MUTED, PP_ALIGN.LEFT)],
        )

    # ——— 5. Architecture ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.5,
        [("Architecture & delivery", 30, True, INK, PP_ALIGN.LEFT)],
    )

    layers = [
        ("Client", "Next.js App Router UI · React Server/Client components · Capacitor WebView"),
        ("API", "Route handlers under /api/* · session-gated · Zod-validated payloads"),
        ("Auth", "Auth.js v5 · credentials + Google · household membership checks"),
        ("Data", "Prisma models · Postgres on Neon · migrations in CI/deploy"),
        ("Jobs", "Alert pipeline (due tickets/events/meds) · push + SMTP email"),
        ("Ops", "Vercel production · env secrets · dual aliases (banik / rose)"),
    ]
    for i, (t, d) in enumerate(layers):
        top = 1.15 + i * 0.95
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(top), Inches(2.2), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL if i % 2 == 0 else TEAL_DARK
        shape.line.fill.background()
        add_textbox(
            s, 0.7, top + 0.2, 1.9, 0.45,
            [(t, 16, True, WHITE, PP_ALIGN.CENTER)],
        )
        card(s, 3.0, top, 9.7, 0.8)
        add_textbox(
            s, 3.25, top + 0.22, 9.2, 0.45,
            [(d, 14, False, INK, PP_ALIGN.LEFT)],
        )

    # ——— 6. CV bullets ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.5,
        [("CV-ready accomplishment bullets", 28, True, INK, PP_ALIGN.LEFT)],
    )
    add_textbox(
        s, 0.6, 0.9, 12, 0.4,
        [
            (
                "Copy-paste into your résumé (tailor numbers / emphasis as needed).",
                13,
                False,
                MUTED,
                PP_ALIGN.LEFT,
            )
        ],
    )

    bullets = [
        "Designed and shipped Flowdesk, a full-stack family CRM covering tasks, calendar, finance, medication schedules, notes, and alerts.",
        "Built a Next.js 15 / TypeScript / Prisma application with 40+ authenticated API routes and household-scoped multi-user data isolation.",
        "Implemented Auth.js authentication (email/password + Google OAuth) and invite-code household onboarding.",
        "Delivered production UX: Today radar, Kanban board, smart natural-language composer, command palette, and theme/wallpaper system.",
        "Created domain modules including BD income-tax calculator, timezone-aware med reminders, recurring events, and a poems atelier with dictionary + doodle tools.",
        "Integrated Web Push and email alert dispatch for due tickets, events, and medication doses.",
        "Packaged a Capacitor Android APK that loads the live Vercel deployment; configured PWA (Serwist) for installable web use.",
        "Operated continuous production deploys on Vercel with Neon Postgres migrations and dual production aliases.",
    ]
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        r = p.add_run()
        set_run(r, "▸  " + b, size=13, color=INK)

    # ——— 7. Tech stack table-ish ———
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_accent_bar(s)
    add_textbox(
        s, 0.6, 0.35, 12, 0.5,
        [("Technology stack", 30, True, INK, PP_ALIGN.LEFT)],
    )

    stack = [
        ("Language", "TypeScript, JavaScript (Node 22)"),
        ("Web framework", "Next.js 15 (App Router), React 19"),
        ("Styling / UX", "Tailwind CSS 4, Framer Motion, Lucide icons"),
        ("ORM / DB", "Prisma 6, PostgreSQL (Neon)"),
        ("Auth", "Auth.js v5, bcrypt, Google OAuth"),
        ("Editors", "Tiptap (notes), custom poem + canvas doodle"),
        ("Mobile", "Capacitor 8 (Android), Geolocation, Local Notifications"),
        ("PWA", "Serwist service worker"),
        ("Integrations", "Open-Meteo weather, Free Dictionary API, Bing wallpaper, Web Push, Nodemailer"),
        ("Tooling", "Zod, date-fns, rrule, chrono-node, ESLint, Vercel CI/CD"),
    ]
    for i, (k, v) in enumerate(stack):
        top = 1.05 + i * 0.58
        kbox = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(top), Inches(2.6), Inches(0.5)
        )
        kbox.fill.solid()
        kbox.fill.fore_color.rgb = TEAL_DARK
        kbox.line.fill.background()
        add_textbox(
            s, 0.65, top + 0.08, 2.4, 0.35,
            [(k, 12, True, WHITE, PP_ALIGN.LEFT)],
        )
        add_textbox(
            s, 3.4, top + 0.08, 9.3, 0.35,
            [(v, 13, False, INK, PP_ALIGN.LEFT)],
        )

    # ——— 8. Closing / links ———
    s = prs.slides.add_slide(blank)
    add_bg(s, TEAL_DARK)
    add_textbox(
        s, 0.9, 2.0, 11.5, 1.0,
        [("Thanks — happy to walk through the codebase", 28, True, WHITE, PP_ALIGN.LEFT)],
    )
    add_textbox(
        s, 0.9, 3.2, 11.5, 2.0,
        [
            ("Live demo", 14, True, LIGHT_TEAL, PP_ALIGN.LEFT),
            ("https://flowdesk-banik.vercel.app", 18, False, WHITE, PP_ALIGN.LEFT),
            ("", 10, False, WHITE, PP_ALIGN.LEFT),
            ("Stack highlight", 14, True, LIGHT_TEAL, PP_ALIGN.LEFT),
            ("Next.js · TypeScript · Prisma · Neon · Auth.js · Capacitor · Vercel", 16, False, WHITE, PP_ALIGN.LEFT),
            ("", 10, False, WHITE, PP_ALIGN.LEFT),
            ("Role", 14, True, LIGHT_TEAL, PP_ALIGN.LEFT),
            ("Sole builder — product design, implementation, and production operations", 16, False, WHITE, PP_ALIGN.LEFT),
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
